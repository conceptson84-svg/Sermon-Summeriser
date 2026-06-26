"""End-of-service export: build a .pptx from the deck, then convert to PDF.

Only runs once when the volunteer clicks "End Service" (eng review issue #1
keeps python-pptx off the live hot path). The .pptx is the shareable artifact;
the PDF is produced via LibreOffice headless if available, else the .pptx is
left as-is with a note.
"""

from __future__ import annotations

import logging
import shutil
import subprocess
from pathlib import Path

log = logging.getLogger(__name__)


def build_pptx(deck, out_path: str | Path, church_name: str = "") -> Path:
    """Write the full deck to a .pptx file. Returns the path."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    out_path = Path(out_path)
    prs = Presentation()
    blank = prs.slide_layouts[6]  # truly blank layout

    for slide in deck.slides:
        s = prs.slides.add_slide(blank)
        top = Inches(0.6)
        if church_name:
            box = s.shapes.add_textbox(Inches(0.6), Inches(0.2), Inches(9), Inches(0.5))
            box.text_frame.text = church_name
            box.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
        body = s.shapes.add_textbox(Inches(0.6), top, Inches(9), Inches(6))
        tf = body.text_frame
        tf.word_wrap = True
        for i, pt in enumerate(slide.points):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = f"• {pt.text}"
            para.runs[0].font.size = Pt(28)
            if pt.scripture:
                ref = tf.add_paragraph()
                ref.text = f"    {pt.scripture}"
                ref.runs[0].font.size = Pt(18)
                ref.runs[0].font.italic = True

    prs.save(str(out_path))
    return out_path


def _find_soffice() -> str | None:
    """Locate the LibreOffice binary across macOS, Windows, and Linux."""
    # On PATH (Linux, or anyone who added it).
    found = shutil.which("soffice") or shutil.which("libreoffice")
    if found:
        return found
    # Common fixed install locations not usually on PATH.
    candidates = [
        # Windows
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        # macOS
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
    ]
    for c in candidates:
        if Path(c).exists():
            return c
    return None


def convert_to_pdf(pptx_path: str | Path) -> Path | None:
    """Convert a .pptx to PDF via LibreOffice headless. Returns the PDF path,
    or None if LibreOffice isn't installed (the .pptx still exists)."""
    pptx_path = Path(pptx_path)
    soffice = _find_soffice()
    if not soffice:
        log.warning("LibreOffice not found — skipping PDF, .pptx is at %s", pptx_path)
        return None
    try:
        subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir",
             str(pptx_path.parent), str(pptx_path)],
            check=True, capture_output=True, timeout=120,
        )
    except (subprocess.CalledProcessError, subprocess.TimeoutExpired) as e:
        log.warning("PDF conversion failed: %s", e)
        return None
    pdf = pptx_path.with_suffix(".pdf")
    return pdf if pdf.exists() else None


def build_pdf_direct(deck, out_path: str | Path, church_name: str = "") -> Path | None:
    """Generate a clean, printable PDF directly with reportlab — no LibreOffice
    needed. Returns the PDF path, or None if reportlab isn't installed."""
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.lib import colors
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    except ImportError:
        log.info("reportlab not installed — cannot build PDF directly")
        return None

    from datetime import datetime

    out_path = Path(out_path)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("Title2", parent=styles["Title"], fontSize=26,
                                 textColor=colors.HexColor("#1c2f73"), spaceAfter=4)
    sub_style = ParagraphStyle("Sub", parent=styles["Normal"], fontSize=12,
                               textColor=colors.HexColor("#6b7591"), spaceAfter=18)
    point_style = ParagraphStyle("Point", parent=styles["Normal"], fontSize=16,
                                 leading=22, spaceAfter=12, leftIndent=10)

    doc = SimpleDocTemplate(str(out_path), pagesize=letter,
                            topMargin=0.8 * inch, bottomMargin=0.8 * inch,
                            leftMargin=0.9 * inch, rightMargin=0.9 * inch)
    story = []
    if church_name:
        story.append(Paragraph(church_name, title_style))
    story.append(Paragraph("Sermon Summary — " + datetime.now().strftime("%A, %d %B %Y"),
                           sub_style))
    points = deck.all_points()
    if not points:
        story.append(Paragraph("(No points captured.)", point_style))
    for pt in points:
        text = f"&bull;&nbsp;&nbsp;{pt.text}"
        if pt.scripture:
            text += (f'  <font color="#3a5bd6" size="12"><i>({pt.scripture})</i></font>')
        story.append(Paragraph(text, point_style))
    story.append(Spacer(1, 0.3 * inch))
    if church_name:
        story.append(Paragraph(f'<font color="#9aa3bd" size="10">Shared from {church_name}</font>',
                               styles["Normal"]))
    doc.build(story)
    return out_path if out_path.exists() else None


def build_text_summary(deck, church_name: str = "") -> str:
    """A clean, shareable plain-text summary (for email / WhatsApp / .txt)."""
    from datetime import datetime

    lines = []
    if church_name:
        lines.append(church_name)
    lines.append("Sermon Summary — " + datetime.now().strftime("%A, %d %B %Y"))
    lines.append("")
    points = deck.all_points()
    if not points:
        lines.append("(No points captured yet.)")
    for pt in points:
        bullet = f"• {pt.text}"
        if pt.scripture:
            bullet += f" ({pt.scripture})"
        lines.append(bullet)
    lines.append("")
    lines.append("Shared from " + (church_name or "our church"))
    return "\n".join(lines)


def export_service(deck, out_dir: str | Path, church_name: str = "") -> dict:
    """Full end-of-service export. Returns {'pptx': path, 'pdf': path|None}."""
    from datetime import datetime

    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    stamp = datetime.now().strftime("%Y%m%d-%H%M")
    pptx = build_pptx(deck, out_dir / f"sermon-summary-{stamp}.pptx", church_name)
    # Prefer the dependency-free reportlab PDF; fall back to LibreOffice only if
    # reportlab isn't installed.
    pdf = build_pdf_direct(deck, out_dir / f"sermon-summary-{stamp}.pdf", church_name)
    if pdf is None:
        pdf = convert_to_pdf(pptx)
    return {"pptx": pptx, "pdf": pdf}
